#!/usr/bin/env python3
"""
Mini-curso de Segurança da Informação · IF Sul · Técnico em Informática
Autores: Daniel Tomm · Francesco Galvão
"""
import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Inches, Emu

ROOT = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(ROOT, "images")
# Fallback if running from /tmp during dev
if not os.path.isdir(IMG_DIR):
    IMG_DIR = "/home/danieltomm/Projects/IFF/cybersec_course/scripts/images"
TEMPLATE = "/home/danieltomm/Projects/IFF/cybersec_course/Cybersec_course.pptx"
OUTPUT = "/home/danieltomm/Projects/IFF/cybersec_course/Aula-Seguranca-IFSul.pptx"


def img(name):
    return os.path.join(IMG_DIR, name)


# ============================================================
# Paleta
# ============================================================
BG = RGBColor(0x0A, 0x0E, 0x1A)
PANEL = RGBColor(0x14, 0x1B, 0x2D)
PANEL_HI = RGBColor(0x1B, 0x24, 0x3A)
BORDER = RGBColor(0x2A, 0x37, 0x55)
BORDER_HI = RGBColor(0x3D, 0x4D, 0x73)

WHITE = RGBColor(0xF8, 0xFA, 0xFC)
TEXT = RGBColor(0xCB, 0xD5, 0xE1)
MUTED = RGBColor(0x8B, 0x97, 0xAE)
DIM = RGBColor(0x4A, 0x56, 0x70)

C_BLUE = RGBColor(0x60, 0xA5, 0xFA)
C_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
C_AMBER = RGBColor(0xFB, 0xBF, 0x24)
C_GREEN = RGBColor(0x4A, 0xDE, 0x80)
C_RED = RGBColor(0xF8, 0x71, 0x71)
C_CYAN = RGBColor(0x22, 0xD3, 0xEE)
C_PINK = RGBColor(0xF4, 0x72, 0xB6)

TERM_GREEN = RGBColor(0x4A, 0xDE, 0x80)
TERM_DIM = RGBColor(0x6B, 0x77, 0x90)

# ============================================================
prs = Presentation(TEMPLATE)
SW, SH = prs.slide_width, prs.slide_height


def _wipe(p):
    sld_lst = p.slides._sldIdLst
    rids = [
        s.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        for s in list(sld_lst)
    ]
    for s in list(sld_lst):
        sld_lst.remove(s)
    for r in rids:
        if not r:
            continue
        try:
            target = p.part.rels[r].target_part
            p.part.drop_rel(r)
            del target.package.parts[target.partname]
        except Exception:
            pass


_wipe(prs)

SECTIONS = {
    1: ("01", "Segurança da Informação", C_BLUE, "sec01-security.jpg"),
    2: ("02", "Como a Internet funciona", C_PURPLE, "sec02-network.jpg"),
    3: ("03", "Boas práticas de defesa", C_AMBER, "sec03-shield.jpg"),
    4: ("04", "Ferramentas", C_GREEN, "sec04-tools.jpg"),
    5: ("05", "Ética e Lei", C_RED, "sec05-law.jpg"),
}

_page_n = [0]


def _next_page():
    _page_n[0] += 1
    return _page_n[0]


# ============================================================
# Helpers de desenho
# ============================================================
def _blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s


def _txt(slide, left, top, width, height, text, *,
         size=18, bold=False, italic=False, color=TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Inter"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return box


def _runs(slide, left, top, width, height, runs_spec, *,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for spec in runs_spec:
        r = p.add_run()
        r.text = spec["text"]
        r.font.size = Pt(spec.get("size", 18))
        r.font.bold = spec.get("bold", False)
        r.font.italic = spec.get("italic", False)
        r.font.color.rgb = spec.get("color", TEXT)
        r.font.name = spec.get("font", "Inter")
    return box


def _paragraphs(slide, left, top, width, height, paragraphs, *,
                anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, runs_spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        for spec in runs_spec:
            r = p.add_run()
            r.text = spec.get("text", "")
            r.font.size = Pt(spec.get("size", 16))
            r.font.bold = spec.get("bold", False)
            r.font.italic = spec.get("italic", False)
            r.font.color.rgb = spec.get("color", TEXT)
            r.font.name = spec.get("font", "Inter")
        if "space_after" in (paragraphs[i][0] if paragraphs[i] else {}):
            p.space_after = Pt(paragraphs[i][0]["space_after"])
        else:
            p.space_after = Pt(6)
    return box


def _rect(slide, left, top, width, height, *,
          fill=PANEL, border=None, border_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_w)
    shp.shadow.inherit = False
    return shp


def _round_rect(slide, left, top, width, height, *,
                fill=PANEL, border=None, border_w=0.75, radius=0.05):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_w)
    shp.shadow.inherit = False
    return shp


def _bar(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _arrow_right(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


# Image with semi-transparent dark overlay so text remains readable on top.
# We approximate transparency by stacking a tinted rectangle over the image.
# python-pptx doesn't expose alpha cleanly via the public API for shapes,
# so we use a fully-opaque dark rect; the image still adds tonal variety
# and the rect prevents it from competing with text.
def _image(slide, path, left, top, width, height, *, dim_alpha=None,
           tint=None, tint_alpha=None):
    pic = slide.shapes.add_picture(path, left, top, width, height)
    if dim_alpha is not None:
        # Stack a near-black overlay
        overlay = _rect(slide, left, top, width, height,
                        fill=tint or BG, border=None)
        # Manually set transparency via XML (a:srgbClr/a:alpha)
        from pptx.oxml.ns import qn
        from lxml import etree
        sp = overlay._element
        spPr = sp.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        if solidFill is not None:
            srgbClr = solidFill.find(qn("a:srgbClr"))
            if srgbClr is not None:
                alpha = etree.SubElement(srgbClr, qn("a:alpha"))
                alpha.set("val", str(int(dim_alpha * 100000)))
    return pic


# ============================================================
# Estruturas reusáveis
# ============================================================
def chrome(slide, section, *, hide_footer=False):
    code, name, color, _ = SECTIONS[section]
    _bar(slide, Emu(0), Emu(0), Inches(0.13), SH, color)
    _runs(slide, Inches(0.55), Inches(0.45), Inches(15), Inches(0.4), [
        {"text": code, "size": 13, "bold": True, "color": color},
        {"text": "  ·  ", "size": 13, "color": DIM},
        {"text": name.upper(), "size": 13, "bold": True, "color": MUTED},
    ])
    if not hide_footer:
        page = _next_page()
        _txt(slide, Inches(0.55), SH - Inches(0.45), Inches(12), Inches(0.35),
             "Daniel Tomm  ·  Francesco Galvão  ·  IF Sul · Técnico em Informática",
             size=12, color=DIM)
        _txt(slide, SW - Inches(1.2), SH - Inches(0.45), Inches(0.6), Inches(0.35),
             f"{page:02d}", size=12, color=DIM, align=PP_ALIGN.RIGHT)


def big_title(slide, title, subtitle="", top=Inches(1.05)):
    _txt(slide, Inches(0.55), top, SW - Inches(1.1), Inches(1.6),
         title, size=52, bold=True, color=WHITE)
    if subtitle:
        _txt(slide, Inches(0.55), top + Inches(1.35), SW - Inches(1.1), Inches(0.7),
             subtitle, size=26, color=MUTED)


# ============================================================
# Tipos de slide
# ============================================================
def cover_slide():
    s = _blank()
    full_w = SW - Inches(1.1)

    # Vertical accent strip on the left edge
    _bar(s, Emu(0), Emu(0), Inches(0.18), SH, C_BLUE)

    # Eyebrow
    _runs(s, Inches(0.55), Inches(1.4), full_w, Inches(0.5), [
        {"text": "MINI-CURSO", "size": 16, "bold": True, "color": C_BLUE},
        {"text": "  ·  ", "size": 16, "color": DIM},
        {"text": "IF SUL · TÉCNICO EM INFORMÁTICA",
         "size": 16, "bold": True, "color": MUTED},
    ])

    # Title — generous, centered around middle of slide
    _txt(s, Inches(0.55), Inches(2.7), full_w, Inches(1.2),
         "Introdução à",
         size=44, bold=True, color=MUTED)
    _txt(s, Inches(0.55), Inches(3.5), full_w, Inches(1.7),
         "Segurança da Informação",
         size=86, bold=True, color=WHITE)
    _txt(s, Inches(0.55), Inches(5.3), full_w, Inches(0.8),
         "Protocolos para defesa, melhores práticas e ferramentas",
         size=26, color=TEXT)
    _txt(s, Inches(0.55), Inches(5.95), full_w, Inches(0.8),
         "para penetração de sistemas.",
         size=26, color=C_BLUE, italic=True)

    _bar(s, Inches(0.6), Inches(7.6), Inches(2.5), Inches(0.05), C_BLUE)

    _txt(s, Inches(0.55), Inches(7.8), full_w, Inches(0.5),
         "AUTORES", size=14, bold=True, color=MUTED)
    _txt(s, Inches(0.55), Inches(8.3), full_w, Inches(0.9),
         "Daniel Tomm   ·   Francesco Galvão",
         size=36, bold=True, color=WHITE)

    _txt(s, Inches(0.55), Inches(9.3), full_w, Inches(0.5),
         "PampaComputing  ·  IF Sul",
         size=14, color=DIM)


def section_divider(num, notes=""):
    s = _blank()
    code, name, color, _image_name = SECTIONS[num]

    # Huge faint number anchoring the layout
    _txt(s, Inches(0.55), Inches(0.8), Inches(8), Inches(7),
         code, size=400, bold=True, color=PANEL_HI, font="Inter")

    _bar(s, Inches(0.0), Inches(0.0), Inches(0.13), SH, color)

    _txt(s, Inches(0.55), Inches(4.3), Inches(15), Inches(0.5),
         "BLOCO  ·  " + code, size=16, bold=True, color=color)

    _txt(s, Inches(0.55), Inches(4.95), Inches(15), Inches(2.5),
         name, size=70, bold=True, color=WHITE)

    # Subtle accent bar under the title
    _bar(s, Inches(0.55), Inches(6.5), Inches(2.0), Inches(0.06), color)

    chrome(s, num, hide_footer=True)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    _next_page()


def content_slide(section, title, subtitle="", body=None, notes=""):
    s = _blank()
    chrome(s, section)
    big_title(s, title, subtitle)

    if body:
        _, _, color, _ = SECTIONS[section]
        _bar(s, Inches(0.55), Inches(3.1), Inches(0.9), Inches(0.04), color)

        paras = []
        for item in body:
            text, level = item if isinstance(item, tuple) else (item, 0)
            if not text:
                paras.append([{"text": " ", "size": 12, "space_after": 8}])
                continue
            if level == 0:
                paras.append([
                    {"text": "▪  ", "size": 28, "color": color, "bold": True},
                    {"text": text, "size": 28, "color": WHITE, "space_after": 14},
                ])
            elif level == 1:
                paras.append([
                    {"text": "      —  ", "size": 24, "color": DIM},
                    {"text": text, "size": 24, "color": TEXT, "space_after": 10},
                ])
            elif level == 2:
                paras.append([
                    {"text": "", "size": 12, "space_after": 8}
                ])

        _paragraphs(s, Inches(0.55), Inches(3.4), SW - Inches(1.1),
                    SH - Inches(4.4), paras)

    if notes:
        s.notes_slide.notes_text_frame.text = notes


def content_slide_with_image(section, title, subtitle, body, image_name,
                             notes=""):
    """Content + side image (right column)."""
    s = _blank()
    chrome(s, section)
    big_title(s, title, subtitle)
    _, _, color, _ = SECTIONS[section]
    _bar(s, Inches(0.55), Inches(3.1), Inches(0.9), Inches(0.04), color)

    text_w = Inches(11.5)
    img_left = text_w + Inches(0.7)
    img_top = Inches(3.4)
    img_w = SW - img_left - Inches(0.55)
    img_h = SH - Inches(4.4)

    # image card
    _round_rect(s, img_left, img_top, img_w, img_h,
                fill=PANEL, border=BORDER, border_w=0.75, radius=0.03)
    _image(s, img(image_name),
           img_left + Inches(0.05), img_top + Inches(0.05),
           img_w - Inches(0.1), img_h - Inches(0.1))
    # Small dim overlay for cohesion with theme
    _image(s, img(image_name),
           img_left + Inches(0.05), img_top + Inches(0.05),
           img_w - Inches(0.1), img_h - Inches(0.1), dim_alpha=0.25)

    paras = []
    for item in body:
        text, level = item if isinstance(item, tuple) else (item, 0)
        if not text:
            paras.append([{"text": " ", "size": 12, "space_after": 8}])
            continue
        if level == 0:
            paras.append([
                {"text": "▪  ", "size": 26, "color": color, "bold": True},
                {"text": text, "size": 26, "color": WHITE, "space_after": 13},
            ])
        elif level == 1:
            paras.append([
                {"text": "      —  ", "size": 22, "color": DIM},
                {"text": text, "size": 22, "color": TEXT, "space_after": 9},
            ])
    _paragraphs(s, Inches(0.55), Inches(3.4), text_w - Inches(0.55),
                SH - Inches(4.4), paras)

    if notes:
        s.notes_slide.notes_text_frame.text = notes


def card_slide(section, title, subtitle, cards, notes=""):
    s = _blank()
    chrome(s, section)
    big_title(s, title, subtitle)

    n = len(cards)
    margin = Inches(0.55)
    gap = Inches(0.3)
    available = SW - 2 * margin - (n - 1) * gap
    card_w = available / n
    card_top = Inches(3.5)
    card_h = Inches(6.4)

    for i, c in enumerate(cards):
        left = margin + i * (card_w + gap)
        col = c.get("color", C_BLUE)

        _round_rect(s, left, card_top, card_w, card_h,
                    fill=PANEL, border=BORDER, border_w=0.75, radius=0.04)
        _round_rect(s, left, card_top, card_w, Inches(0.18),
                    fill=col, border=None, radius=0.5)

        big = c.get("big", "")
        big_size = 108 if len(big) <= 3 else (84 if len(big) <= 4 else 68)
        _txt(s, left + Inches(0.4), card_top + Inches(0.5),
             card_w - Inches(0.8), Inches(2.4),
             big, size=big_size, bold=True, color=col)

        _txt(s, left + Inches(0.4), card_top + Inches(3.0),
             card_w - Inches(0.8), Inches(0.4),
             c.get("label", "").upper(),
             size=13, bold=True, color=MUTED)

        _txt(s, left + Inches(0.4), card_top + Inches(3.4),
             card_w - Inches(0.8), Inches(0.9),
             c.get("title", ""),
             size=32, bold=True, color=WHITE)

        _txt(s, left + Inches(0.4), card_top + Inches(4.4),
             card_w - Inches(0.8), card_h - Inches(4.6),
             c.get("body", ""),
             size=22, color=TEXT)

    if notes:
        s.notes_slide.notes_text_frame.text = notes


def two_card_slide(section, title, subtitle, left_card, right_card, notes=""):
    s = _blank()
    chrome(s, section)
    big_title(s, title, subtitle)

    margin = Inches(0.55)
    gap = Inches(0.4)
    card_w = (SW - 2 * margin - gap) / 2
    card_top = Inches(3.4)
    card_h = SH - Inches(4.4)

    for i, c in enumerate([left_card, right_card]):
        left = margin + i * (card_w + gap)
        col = c.get("color", C_BLUE)

        _round_rect(s, left, card_top, card_w, card_h,
                    fill=PANEL, border=BORDER, border_w=0.75, radius=0.03)
        _round_rect(s, left, card_top, card_w, Inches(0.16),
                    fill=col, border=None, radius=0.5)

        _txt(s, left + Inches(0.5), card_top + Inches(0.5),
             card_w - Inches(1), Inches(0.4),
             c.get("label", "").upper(),
             size=14, bold=True, color=col)
        _txt(s, left + Inches(0.5), card_top + Inches(1.0),
             card_w - Inches(1), Inches(1.0),
             c.get("title", ""),
             size=36, bold=True, color=WHITE)

        items = c.get("items", [])
        paras = []
        for item in items:
            if not item:
                paras.append([{"text": " ", "size": 14, "space_after": 10}])
                continue
            paras.append([
                {"text": "▪  ", "size": 26, "color": col, "bold": True},
                {"text": item, "size": 26, "color": TEXT, "space_after": 12},
            ])
        _paragraphs(s, left + Inches(0.5), card_top + Inches(2.4),
                    card_w - Inches(1), card_h - Inches(2.6),
                    paras)

    if notes:
        s.notes_slide.notes_text_frame.text = notes


def authors_slide(notes=""):
    """Two-column bio slide with company badges and roles."""
    s = _blank()
    chrome(s, 1)
    big_title(s,
        "Quem está com vocês hoje",
        "Dois profissionais de áreas bem diferentes da segurança.")
    _bar(s, Inches(0.55), Inches(3.1), Inches(0.9), Inches(0.04), C_BLUE)

    margin = Inches(0.55)
    gap = Inches(0.4)
    card_w = (SW - 2 * margin - gap) / 2
    card_top = Inches(3.4)
    card_h = SH - Inches(4.4)

    authors = [
        {
            "name": "Daniel Tomm",
            "company": "NewSeg",
            "company_full": "soluções para segurança pública",
            "role": "Engenharia de Software",
            "color": C_BLUE,
            "initials": "DT",
            "bio": [
                "Trabalha desenvolvendo sistemas usados por forças de segurança pública no Brasil",
                "Foco em backend, integrações e arquitetura de sistemas críticos",
                "Atua na ponta entre desenvolvimento e operação — onde os ataques chegam primeiro",
            ],
        },
        {
            "name": "Francesco Galvão",
            "company": "Red Hat",
            "company_full": "uma das maiores empresas de tecnologia open-source do mundo",
            "role": "Support Engineer",
            "color": C_RED,
            "initials": "FG",
            "bio": [
                "Trabalha resolvendo problemas reais de segurança em infraestrutura corporativa",
                "Lida no dia a dia com Linux, contêineres, e o ecossistema enterprise",
                "Foco em descobrir, mitigar e corrigir vulnerabilidades em produção",
            ],
        },
    ]

    for i, a in enumerate(authors):
        left = margin + i * (card_w + gap)
        col = a["color"]

        _round_rect(s, left, card_top, card_w, card_h,
                    fill=PANEL, border=BORDER, border_w=0.75, radius=0.03)
        _round_rect(s, left, card_top, card_w, Inches(0.16),
                    fill=col, border=None, radius=0.5)

        # Initials avatar
        avatar_size = Inches(1.4)
        _round_rect(s, left + Inches(0.6), card_top + Inches(0.7),
                    avatar_size, avatar_size,
                    fill=PANEL_HI, border=col, border_w=2, radius=0.5)
        _txt(s, left + Inches(0.6), card_top + Inches(0.7),
             avatar_size, avatar_size,
             a["initials"], size=42, bold=True, color=col,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Name + role
        text_left = left + Inches(2.2)
        _txt(s, text_left, card_top + Inches(0.75),
             card_w - Inches(2.7), Inches(0.7),
             a["name"], size=32, bold=True, color=WHITE)
        _txt(s, text_left, card_top + Inches(1.4),
             card_w - Inches(2.7), Inches(0.5),
             a["role"], size=17, color=MUTED)

        # Company line
        _runs(s, left + Inches(0.6), card_top + Inches(2.4),
              card_w - Inches(1.2), Inches(0.5), [
            {"text": "TRABALHA NA  ", "size": 13, "bold": True, "color": MUTED},
            {"text": a["company"], "size": 13, "bold": True, "color": col},
        ])
        _txt(s, left + Inches(0.6), card_top + Inches(2.85),
             card_w - Inches(1.2), Inches(0.6),
             a["company_full"], size=17, italic=True, color=TEXT)

        # Bio bullets
        paras = []
        for line in a["bio"]:
            paras.append([
                {"text": "▪  ", "size": 20, "color": col, "bold": True},
                {"text": line, "size": 20, "color": TEXT, "space_after": 10},
            ])
        _paragraphs(s, left + Inches(0.6), card_top + Inches(4.0),
                    card_w - Inches(1.2), card_h - Inches(4.2),
                    paras)

    if notes:
        s.notes_slide.notes_text_frame.text = notes


def code_slide(section, title, subtitle, prompt_lines, notes=""):
    s = _blank()
    chrome(s, section)
    big_title(s, title, subtitle)

    term_left = Inches(0.55)
    term_top = Inches(3.4)
    term_w = SW - Inches(1.1)
    term_h = SH - Inches(4.4)

    _round_rect(s, term_left, term_top, term_w, term_h,
                fill=RGBColor(0x05, 0x09, 0x12),
                border=BORDER, border_w=0.75, radius=0.02)

    for i, c in enumerate([RGBColor(0xF8, 0x71, 0x71),
                           RGBColor(0xFB, 0xBF, 0x24),
                           RGBColor(0x4A, 0xDE, 0x80)]):
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 term_left + Inches(0.3) + Inches(0.4 * i),
                                 term_top + Inches(0.3),
                                 Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.fill.background()
        dot.shadow.inherit = False

    _txt(s, term_left, term_top + Inches(0.27),
         term_w, Inches(0.3),
         "terminal — bash",
         size=10, color=TERM_DIM, align=PP_ALIGN.CENTER, font="JetBrains Mono")

    body_left = term_left + Inches(0.35)
    body_top = term_top + Inches(0.85)
    body_w = term_w - Inches(0.7)

    paras = []
    for entry in prompt_lines:
        if "comment" in entry:
            paras.append([
                {"text": "# " + entry["comment"], "size": 20,
                 "color": TERM_DIM, "italic": True,
                 "font": "JetBrains Mono", "space_after": 6},
            ])
        if "cmd" in entry:
            paras.append([
                {"text": entry.get("prompt", "$") + " ", "size": 22,
                 "bold": True, "color": TERM_GREEN, "font": "JetBrains Mono"},
                {"text": entry["cmd"], "size": 22, "color": WHITE,
                 "font": "JetBrains Mono", "space_after": 10},
            ])
        if "output" in entry:
            paras.append([
                {"text": entry["output"], "size": 19,
                 "color": TEXT, "font": "JetBrains Mono", "space_after": 8},
            ])
        if entry.get("blank"):
            paras.append([{"text": " ", "size": 10, "space_after": 5}])

    _paragraphs(s, body_left, body_top, body_w, term_h - Inches(1.0), paras)

    if notes:
        s.notes_slide.notes_text_frame.text = notes


def diagram_browser_burp_server(section, title, subtitle, notes=""):
    s = _blank()
    chrome(s, section)
    big_title(s, title, subtitle)

    boxes_top = Inches(3.9)
    box_h = Inches(3.0)
    box_w = Inches(4.2)
    arrow_w = Inches(0.9)
    total_w = 3 * box_w + 2 * arrow_w + Inches(0.8)
    start = (SW - total_w) / 2 + Inches(0.4)

    _round_rect(s, start, boxes_top, box_w, box_h,
                fill=PANEL, border=C_BLUE, border_w=2, radius=0.04)
    _txt(s, start, boxes_top + Inches(0.45), box_w, Inches(0.7),
         "🌐", size=48, color=C_BLUE, align=PP_ALIGN.CENTER)
    _txt(s, start, boxes_top + Inches(1.35), box_w, Inches(0.4),
         "NAVEGADOR", size=14, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
    _txt(s, start, boxes_top + Inches(1.8), box_w, Inches(0.7),
         "cliente", size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(s, start + Inches(0.3), boxes_top + Inches(2.5),
         box_w - Inches(0.6), Inches(0.5),
         "você está aqui", size=14, italic=True,
         color=MUTED, align=PP_ALIGN.CENTER)

    _arrow_right(s, start + box_w + Inches(0.2),
                 boxes_top + Inches(1.3), arrow_w, Inches(0.45), C_GREEN)

    bx = start + box_w + arrow_w + Inches(0.4)
    _round_rect(s, bx, boxes_top, box_w, box_h,
                fill=PANEL_HI, border=C_GREEN, border_w=2, radius=0.04)
    _txt(s, bx, boxes_top + Inches(0.45), box_w, Inches(0.7),
         "🛡", size=48, color=C_GREEN, align=PP_ALIGN.CENTER)
    _txt(s, bx, boxes_top + Inches(1.35), box_w, Inches(0.4),
         "BURP SUITE", size=14, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    _txt(s, bx, boxes_top + Inches(1.8), box_w, Inches(0.7),
         "proxy", size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(s, bx + Inches(0.3), boxes_top + Inches(2.5),
         box_w - Inches(0.6), Inches(0.5),
         "intercepta + edita", size=14, italic=True,
         color=MUTED, align=PP_ALIGN.CENTER)

    _arrow_right(s, bx + box_w + Inches(0.2),
                 boxes_top + Inches(1.3), arrow_w, Inches(0.45), C_GREEN)

    sx = bx + box_w + arrow_w + Inches(0.4)
    _round_rect(s, sx, boxes_top, box_w, box_h,
                fill=PANEL, border=C_PURPLE, border_w=2, radius=0.04)
    _txt(s, sx, boxes_top + Inches(0.45), box_w, Inches(0.7),
         "🖥", size=48, color=C_PURPLE, align=PP_ALIGN.CENTER)
    _txt(s, sx, boxes_top + Inches(1.35), box_w, Inches(0.4),
         "SERVIDOR", size=14, bold=True, color=C_PURPLE, align=PP_ALIGN.CENTER)
    _txt(s, sx, boxes_top + Inches(1.8), box_w, Inches(0.7),
         "remoto", size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(s, sx + Inches(0.3), boxes_top + Inches(2.5),
         box_w - Inches(0.6), Inches(0.5),
         "responde sem saber", size=14, italic=True,
         color=MUTED, align=PP_ALIGN.CENTER)

    _txt(s, Inches(0.55), boxes_top + box_h + Inches(0.6),
         SW - Inches(1.1), Inches(0.6),
         "Cada requisição passa pelo Burp antes de chegar ao servidor — você pode pausar, editar e reenviar.",
         size=18, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    if notes:
        s.notes_slide.notes_text_frame.text = notes


def quote_slide(section, label, big_text, sub="", notes=""):
    s = _blank()
    chrome(s, section)
    _, _, color, _ = SECTIONS[section]

    # Decorative accent
    _bar(s, Inches(0.55), Inches(2.4), Inches(1.5), Inches(0.06), color)

    _txt(s, Inches(0.55), Inches(2.6), SW - Inches(1.1), Inches(0.5),
         label.upper(), size=16, bold=True, color=color)
    _txt(s, Inches(0.55), Inches(3.2), SW - Inches(1.1), Inches(5),
         "“" + big_text + "”",
         size=80, bold=True, color=WHITE)
    if sub:
        _txt(s, Inches(0.55), Inches(8.5), SW - Inches(1.1), Inches(0.7),
             sub, size=26, italic=True, color=MUTED)

    if notes:
        s.notes_slide.notes_text_frame.text = notes


# ============================================================
# Construção do deck
# ============================================================

# 1
cover_slide()

# 2 — Quem somos
authors_slide(notes=(
    "Daniel apresenta o background dele primeiro, então o Francesco. "
    "Reforce que vocês trazem perspectivas COMPLEMENTARES: produto/dev "
    "(NewSeg) e infra/suporte enterprise (Red Hat). Diga que segurança "
    "não é uma carreira só — tem várias portas de entrada."
))

# 3
content_slide(1,
    "O que vamos fazer hoje",
    "Plano da próxima 1h.",
    body=[
        "Conceitos: o que é segurança da informação, como a internet funciona, boas práticas",
        "Ferramentas: Burp Suite, Nmap e automação com IA (Claude Code CLI)",
        "Atividade prática: vocês vão atacar um banco fictício no meio da aula",
        "Fechamento: ética, leis brasileiras e por onde seguir estudando",
        "",
        "Vocês não precisam saber nada de segurança pra acompanhar.",
        ("Se já programaram alguma coisa, vai ser mais rápido — mas não é requisito.", 1),
    ],
)

# 4 — Section 1
section_divider(1)

# 5 — CIA pillars
card_slide(1,
    "Segurança da Informação tem três pilares",
    "Toda decisão de segurança gira em torno desses três.",
    cards=[
        {
            "big": "C", "color": C_BLUE,
            "label": "Confidentiality",
            "title": "Confidencialidade",
            "body": "Só quem é autorizado vê o dado. Senha não vaza. Mensagem privada continua privada.",
        },
        {
            "big": "I", "color": C_PURPLE,
            "label": "Integrity",
            "title": "Integridade",
            "body": "O dado não é alterado por quem não pode. Ninguém edita seu boletim sem você saber.",
        },
        {
            "big": "A", "color": C_GREEN,
            "label": "Availability",
            "title": "Disponibilidade",
            "body": "O sistema está no ar quando você precisa. Banco não pode sair fora 6h por dia.",
        },
    ],
    notes=(
        "Pergunta pra turma: 'qual desses 3 vocês acham mais importante?' "
        "Resposta: depende. Banco prioriza integridade. Hospital prioriza "
        "disponibilidade. Inteligência militar prioriza confidencialidade."
    ),
)

# 6 — Tipos de atacantes
content_slide_with_image(1,
    "Quem ataca, e por quê?",
    "A imagem do hacker de capuz é só o começo.",
    body=[
        "Cibercriminosos — ganham dinheiro: ransomware, fraude, golpe",
        "Hacktivistas — querem fazer um recado político ou social",
        "Insiders — funcionário descontente, ex-funcionário, terceirizado",
        "Estados / espionagem — guerra digital entre países",
        "Script kiddies — adolescentes copiando ferramentas, fazendo bagunça",
        "Pesquisadores e pentesters — atacam com autorização, com contrato",
    ],
    image_name="attackers.jpg",
    notes=(
        "A imagem do capuz no escuro é o estereótipo da mídia. Reforce que "
        "a maioria dos ataques é pra dinheiro, não pra mostrar habilidade. "
        "E que pesquisador de segurança hoje é uma carreira respeitada."
    ),
)

# 7 — Section 2
section_divider(2)

# 8
two_card_slide(2,
    "Cliente e Servidor",
    "O modelo mental que sustenta tudo no resto da aula.",
    left_card={
        "label": "Cliente",
        "title": "Onde você está",
        "color": C_BLUE,
        "items": [
            "navegador, app de celular, terminal, programa qualquer",
            "mostra a tela e roda do seu lado",
            "PODE ser modificado, falsificado, hackeado",
            "é descartável — nunca é fonte de verdade",
        ],
    },
    right_card={
        "label": "Servidor",
        "title": "Onde os dados moram",
        "color": C_PURPLE,
        "items": [
            "computador remoto que guarda e processa os dados",
            "decide o que pode e o que não pode",
            "é onde MORAM as regras de verdade",
            "é a única coisa em que vale a pena confiar",
        ],
    },
    notes=(
        "Analogia: o caixa do banco é o cliente, o cofre é o servidor."
    ),
)

# 9
content_slide_with_image(2,
    "HTTP — toda conversa na web tem o mesmo formato",
    "Cliente faz uma pergunta. Servidor responde.",
    body=[
        "Uma REQUEST tem: método + URL + headers + (às vezes) body",
        "Uma RESPONSE tem: status code + headers + body",
        "É tudo TEXTO. Você consegue ler com os olhos.",
        "",
        "Exemplo de request:",
        ("GET /api/saldo  HTTP/1.1", 1),
        ("Host: meu-banco.com", 1),
        ("Authorization: Bearer xxx", 1),
        "",
        "Exemplo de response:",
        ("HTTP/1.1 200 OK", 1),
        ('{ \"saldo\": 1234.56 }', 1),
    ],
    image_name="http-conv.jpg",
)

# 10 — Métodos
card_slide(2,
    "Os 4 métodos que importam",
    "Mesma URL com método diferente faz coisa diferente.",
    cards=[
        {"big": "GET", "color": C_BLUE, "label": "Ler", "title": "Consultar",
         "body": "Pede uma informação. Não muda nada no servidor. Ex: ver saldo."},
        {"big": "POST", "color": C_GREEN, "label": "Criar", "title": "Enviar novo",
         "body": "Cria algo novo. Ex: fazer login, abrir conta."},
        {"big": "PATCH", "color": C_AMBER, "label": "Atualizar", "title": "Mudar parte",
         "body": "Atualiza só os campos enviados. Ex: trocar nome do cliente."},
        {"big": "DEL", "color": C_RED, "label": "Apagar", "title": "Remover",
         "body": "Apaga algo. Ex: encerrar conta, remover postagem."},
    ],
)

# 11 — Status codes
content_slide_with_image(2,
    "Status codes que vocês vão ver hoje",
    "Toda response começa com um número de 3 dígitos.",
    body=[
        "200  ·  OK, deu certo",
        "201  ·  Criado (depois de um POST)",
        "401  ·  Você não está logado",
        "403  ·  Você está logado mas NÃO PODE fazer isso",
        "404  ·  Não existe",
        "500  ·  O servidor quebrou",
        "",
        "Foque em 200 / 401 / 403 — são os que mais aparecem no dia a dia.",
    ],
    image_name="status-codes.jpg",
    notes=(
        "Na atividade prática, /api/flag retorna 403 antes da exploração e "
        "200 depois — dica visual gigante."
    ),
)

# 12 — Headers + JSON
two_card_slide(2,
    "Headers e JSON",
    "Headers são o envelope. JSON é o conteúdo.",
    left_card={
        "label": "Headers — o envelope", "title": "Metadados da requisição",
        "color": C_AMBER,
        "items": [
            "Content-Type — o que tem dentro do body",
            "User-Agent — qual navegador/app está mandando",
            "Accept-Language — idioma preferido",
            "Cookie / Authorization — quem é você",
        ],
    },
    right_card={
        "label": "JSON — o conteúdo", "title": "Estrutura de dados",
        "color": C_GREEN,
        "items": [
            "estrutura igual a objetos de JS / dicts de Python",
            "chave: valor — pode aninhar e fazer listas",
            '{ "nome": "Ana", "isAdmin": false }',
            "formato universal de troca de dados",
        ],
    },
)

# 13 — Section 3
section_divider(3)

# 14 — Quote: never trust
quote_slide(3,
    "regra de ouro",
    "Never trust the client.",
    sub="Toda regra de negócio mora no servidor. Sempre.",
    notes=(
        "Escreva 'never trust the client' no quadro e deixe lá pelo resto "
        "da aula. Pergunte: 'se o checkbox \"sou admin\" fosse só do "
        "navegador, o que dava errado?'"
    ),
)

# 15 — Onde os servidores falham
content_slide_with_image(3,
    "Onde os servidores falham",
    "As armadilhas mais comuns que dev iniciante cai.",
    body=[
        "Aceitam dados sem validar tipo, formato ou tamanho",
        "Aceitam campos que o cliente não deveria poder mexer",
        ('"ah, o front nunca manda isso" — manda sim, com Burp manda', 1),
        "Verificam permissão SÓ no JavaScript do navegador",
        "Cookies sem proteção (HttpOnly, SameSite, Secure)",
        "Mensagens de erro que vazam estrutura interna do banco",
        "Confiam em IDs/parâmetros vindos da URL sem validar",
    ],
    image_name="fails.jpg",
    notes=(
        "NÃO cite o termo técnico 'mass assignment' aqui — esse é o nome "
        "do bug específico do CTF."
    ),
)

# 16 — OWASP
content_slide_with_image(3,
    "Os 4 vilões mais comuns na web hoje",
    "OWASP Top 10 — versão simplificada.",
    body=[
        "Broken Access Control  ·  servidor não checa permissão",
        ("#1 do mundo real. É o que vocês vão explorar na prática.", 1),
        "Broken Authentication  ·  senha fraca, sessão mal feita",
        "Injection (SQL, XSS, …)  ·  cliente envia código que vira execução",
        "Security Misconfiguration  ·  defaults inseguros, debug ligado",
    ],
    image_name="owasp.jpg",
)

# 17 — Defesa em camadas
content_slide_with_image(3,
    "Defesa em camadas — defense in depth",
    "Você nunca confia em uma única barreira.",
    body=[
        "Validação no front-end  ·  experiência do usuário, NÃO é segurança",
        "Validação no back-end  ·  primeira barreira de verdade",
        "Permissões por rota  ·  só admin acessa /admin, ponto",
        "Banco de dados com privilégios mínimos  ·  app de leitura não escreve",
        "Logs e monitoramento  ·  se algo passar, você descobre depois",
        "Atualizações em dia  ·  90% dos hacks usam falhas já corrigidas",
    ],
    image_name="defense.jpg",
)

# 18 — Section 4
section_divider(4)

# 19 — Burp diagram
diagram_browser_burp_server(4,
    "Ferramenta 1  ·  Burp Suite",
    "Um proxy que se senta entre você e o servidor.",
)

# 20 — Burp features
content_slide_with_image(4,
    "O que o Burp deixa você fazer",
    "As 4 funções que vocês vão usar hoje.",
    body=[
        "Intercept  ·  pausa a requisição antes dela sair do navegador",
        "Forward  ·  libera a requisição (com ou sem edição)",
        "Repeater  ·  manda a mesma requisição várias vezes alterando o body",
        "Proxy History  ·  log completo de tudo que passou na sessão",
        "",
        "Por que profissionais usam: o navegador esconde coisas. O Burp não.",
    ],
    image_name="burp-features.jpg",
)

# 21 — Atividade prática (com imagem)
content_slide_with_image(4,
    "Hora da prática  ·  PampaBank CTF",
    "Vocês vão usar o Burp pra atacar um banco fictício.",
    body=[
        "URL  ·  pampabank-ctf.vercel.app",
        "Credenciais  ·  aluno@pampabank.ctf  /  ctf123",
        "Missão  ·  virar cliente PREMIUM sem pagar nada",
        "Ferramenta  ·  Burp Suite (DevTools também serve)",
        "",
        '💡  Botão "Estou empacado" no canto inferior direito',
        ("4 dicas progressivas se travarem", 1),
        "Não dêem a resposta uns pros outros.",
    ],
    image_name="activity-bank.jpg",
    notes=(
        "Reserve 15-20min pra atividade. Circule pela sala observando "
        "qual hint cada um abriu."
    ),
)

# 22 — Nmap intro (com imagem)
content_slide_with_image(4,
    "Ferramenta 2  ·  Nmap",
    "Mapeador de rede — descobre o que tá rodando onde.",
    body=[
        "Pentester abre o Nmap antes de qualquer outra coisa",
        "Mostra:",
        ("quais máquinas estão ligadas e respondendo", 1),
        ("quais portas estão abertas em cada uma", 1),
        ("qual serviço (e qual versão) tá rodando", 1),
        ("se dá pra deduzir o sistema operacional", 1),
        "",
        "Não invade nada — só pergunta. Mas ouvir bem já dá muita informação.",
    ],
    image_name="nmap-network.jpg",
)

# 23 — Nmap commands (terminal)
code_slide(4,
    "Nmap  ·  comandos básicos pra começar",
    "Você vai usar 90% do tempo só essas linhas.",
    prompt_lines=[
        {"comment": "descobrir quais máquinas estão vivas na rede local"},
        {"cmd": "nmap -sn 192.168.0.0/24"},
        {"blank": True},
        {"comment": "ver as portas abertas de uma máquina específica"},
        {"cmd": "nmap 192.168.0.10"},
        {"blank": True},
        {"comment": "descobrir o serviço e versão de cada porta"},
        {"cmd": "nmap -sV 192.168.0.10"},
        {"blank": True},
        {"comment": "tudo: portas + serviços + SO + scripts default"},
        {"cmd": "nmap -A 192.168.0.10"},
        {"blank": True},
        {"comment": "varredura mais discreta (não bate de frente)"},
        {"cmd": "nmap -sS -T2 192.168.0.10"},
    ],
)

# 24 — IA + Claude (com imagem)
content_slide_with_image(4,
    "Ferramenta 3  ·  IA acelerando tudo",
    "Modelos atuais não substituem você — multiplicam.",
    body=[
        "ChatGPT, Claude, Copilot e afins viraram colegas de trabalho",
        "Em segurança, eles servem pra:",
        ("explicar uma vulnerabilidade que você nunca viu", 1),
        ("gerar scripts de pentest sob medida", 1),
        ("revisar código procurando falhas comuns", 1),
        ("traduzir Bash antigo pra Python (ou vice-versa)", 1),
        ("automatizar tarefa repetitiva", 1),
        "",
        "Claude Code CLI roda no terminal, lê seus arquivos, escreve código.",
    ],
    image_name="ai-future.jpg",
)

# 25 — Claude exemplo
code_slide(4,
    "Exemplo  ·  Claude gerando um scanner em 1 prompt",
    "Da linguagem natural pro script funcional.",
    prompt_lines=[
        {"comment": "abro o Claude Code direto no terminal do projeto"},
        {"cmd": "claude"},
        {"blank": True},
        {"output": ">  Escreva um script Python que use python-nmap pra"},
        {"output": ">  varrer minha rede local e gerar um relatório em"},
        {"output": ">  Markdown listando hosts vivos + portas abertas."},
        {"blank": True},
        {"comment": "Claude lê o projeto, escreve scan_network.py, te explica e pergunta se pode rodar"},
        {"cmd": "python scan_network.py"},
        {"output": "✓ scan completo · relatório salvo em ./report.md"},
    ],
)

# 26 — Section 5
section_divider(5)

# 27 — Autorização
content_slide(5,
    "Hacking só vale com autorização",
    "Sem permissão = crime. Não tem zona cinzenta.",
    body=[
        "Ambientes que são SEUS  ·  máquina virtual, servidor próprio",
        "CTFs e laboratórios  ·  TryHackMe, Hack The Box, PortSwigger Academy",
        "Bug bounty  ·  HackerOne, BugCrowd, Intigriti  —  você assina termos",
        "Vagas de pentest  ·  contrato te dando escopo claro",
        "",
        "Tudo o resto = crime, mesmo que pareça inofensivo.",
    ],
)

# 28 — Leis brasileiras
content_slide(5,
    "Leis brasileiras que vocês precisam conhecer",
    "Carreira em segurança não começa com cadeia.",
    body=[
        'Lei 12.737/2012  ·  "Carolina Dieckmann"',
        ("invasão de dispositivo informático sem autorização", 1),
        ("pena: 3 meses a 1 ano de detenção, mais multa", 1),
        "LGPD  ·  Lei 13.709/2018",
        ("vazamento ou uso indevido de dados pessoais", 1),
        ("multa: até R$ 50 milhões pra empresa, mais responsabilização pessoal", 1),
        "Marco Civil da Internet  ·  Lei 12.965/2014",
        ("regula direitos e deveres na internet brasileira", 1),
    ],
)

# 29 — Próximos passos (com imagem)
content_slide_with_image(5,
    "Por onde seguir agora",
    "A aula acaba aqui — o aprendizado é vocês.",
    body=[
        "Plataformas grátis pra praticar:",
        ("TryHackMe  ·  trilhas guiadas, ótimo pra começar", 1),
        ("PortSwigger Web Security Academy  ·  feito pelos donos do Burp", 1),
        ("Hack The Box  ·  já assume que você sabe um pouco", 1),
        ("OverTheWire  ·  CTFs no terminal, foco em Linux", 1),
        "",
        "Comunidades:",
        ("OWASP, DEF CON groups locais", 1),
        ("CTFs nacionais: H2HC, Roadsec, Mente Binária", 1),
    ],
    image_name="closing.jpg",
)

# 30 — Final
quote_slide(5,
    "valeu pessoal",
    "Saber atacar dá poder.",
    sub="Poder vem com responsabilidade.",
)

# ============================================================
prs.save(OUTPUT)
print(f"OK · {len(prs.slides)} slides → {OUTPUT}")
